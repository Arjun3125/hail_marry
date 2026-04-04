"""Extended data connectors — PPTX, Excel, Google Docs, Notion.

Each connector returns a standardized format:
    {"text": str, "metadata": dict, "chunks": list[str]}
"""
import io
import os
import re
from typing import Any, Optional


def extract_pptx(file_path: str) -> dict[str, Any]:
    """Extract text from a PowerPoint file.

    Requires: pip install python-pptx
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"text": "", "metadata": {"error": "python-pptx not installed"}, "chunks": []}

    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)
        if slide_content:
            slides_text.append(f"[Slide {i}] " + " ".join(slide_content))

    full_text = "\n\n".join(slides_text)
    return {
        "text": full_text,
        "metadata": {"file_type": "pptx", "slide_count": len(prs.slides), "file_path": file_path},
        "chunks": slides_text,
    }


def extract_excel(file_path: str) -> dict[str, Any]:
    """Extract data from an Excel file.

    Requires: pip install openpyxl
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"text": "", "metadata": {"error": "openpyxl not installed"}, "chunks": []}

    wb = load_workbook(file_path, read_only=True, data_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(c for c in cells if c)
            if line.strip():
                rows_text.append(line)
        if rows_text:
            sheet_content = f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text)
            sheets_text.append(sheet_content)

    wb.close()
    full_text = "\n\n".join(sheets_text)
    return {
        "text": full_text,
        "metadata": {"file_type": "xlsx", "sheet_count": len(wb.sheetnames), "file_path": file_path},
        "chunks": sheets_text,
    }


def extract_google_doc(doc_url: str, credentials: Optional[dict] = None) -> dict[str, Any]:
    """Extract content from a Google Doc via export URL.

    For public docs, uses the plain-text export URL.
    For private docs, credentials dict with 'access_token' is required.
    """
    import httpx

    # Extract doc ID from URL
    match = re.search(r"/document/d/([a-zA-Z0-9_-]+)", doc_url)
    if not match:
        return {"text": "", "metadata": {"error": "Invalid Google Docs URL"}, "chunks": []}

    doc_id = match.group(1)
    export_url = f"https://docs.google.com/document/d/{doc_id}/export?format=txt"

    headers = {}
    if credentials and credentials.get("access_token"):
        headers["Authorization"] = f"Bearer {credentials['access_token']}"

    try:
        resp = httpx.get(export_url, headers=headers, timeout=30.0, follow_redirects=True)
        resp.raise_for_status()
        text = resp.text
    except Exception as e:
        return {"text": "", "metadata": {"error": str(e)}, "chunks": []}

    # Split into chunks by double newline
    chunks = [chunk.strip() for chunk in text.split("\n\n") if chunk.strip()]

    return {
        "text": text,
        "metadata": {"file_type": "google_doc", "doc_id": doc_id, "url": doc_url},
        "chunks": chunks,
    }


def extract_notion_page(page_id: str, api_key: str) -> dict[str, Any]:
    """Extract content from a Notion page via API.

    Requires a Notion integration API key.
    """
    import httpx

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Notion-Version": "2022-06-28",
        "Content-Type": "application/json",
    }

    try:
        # Get page blocks
        resp = httpx.get(
            f"https://api.notion.com/v1/blocks/{page_id}/children?page_size=100",
            headers=headers,
            timeout=30.0,
        )
        resp.raise_for_status()
        data = resp.json()
    except Exception as e:
        return {"text": "", "metadata": {"error": str(e)}, "chunks": []}

    blocks = data.get("results", [])
    chunks = []

    for block in blocks:
        block_type = block.get("type", "")
        block_data = block.get(block_type, {})

        # Extract text from rich_text array
        rich_texts = block_data.get("rich_text", [])
        text_parts = [rt.get("plain_text", "") for rt in rich_texts]
        text = "".join(text_parts).strip()

        if text:
            prefix = ""
            if block_type.startswith("heading"):
                level = block_type[-1] if block_type[-1].isdigit() else ""
                prefix = f"{'#' * int(level)} " if level else "# "
            elif block_type == "bulleted_list_item":
                prefix = "• "
            elif block_type == "numbered_list_item":
                prefix = "- "

            chunks.append(f"{prefix}{text}")

    full_text = "\n".join(chunks)
    return {
        "text": full_text,
        "metadata": {"file_type": "notion", "page_id": page_id},
        "chunks": chunks,
    }


# ── Connector registry ──

CONNECTORS = {
    "pptx": extract_pptx,
    "xlsx": extract_excel,
    "google_doc": extract_google_doc,
    "notion": extract_notion_page,
}


def get_connector(file_type: str):
    """Get the appropriate connector function for a file type."""
    return CONNECTORS.get(file_type)
